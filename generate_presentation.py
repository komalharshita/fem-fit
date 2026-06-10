import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize Presentation with 16:9 Widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2. Color Palette Definitions (Strictly based on USER styling image)
COLOR_BLACK = RGBColor(10, 10, 10)       # Solid Deep Black
COLOR_LIME = RGBColor(181, 230, 29)      # High-energy Lime Green
COLOR_YELLOW = RGBColor(242, 201, 76)    # Warm Amber/Gold Yellow
COLOR_WHITE = RGBColor(255, 255, 255)    # High-contrast White
COLOR_MUTED = RGBColor(160, 160, 160)    # Muted Grey for description
COLOR_CARD_BG = RGBColor(24, 24, 24)     # Dark Grey card container

# 3. Helper Functions
def set_slide_background(slide):
    """Fills slide background with deep black."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BLACK

def create_slide_header(slide, title_text, category_text="FEMFIT-INDUSTRIAL"):
    """Creates a consistent premium header with a lime vertical bar and yellow category tag."""
    # Add vertical neon-lime bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.1), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_LIME
    bar.line.fill.background()
    
    # Text Frame for Category and Title
    tx_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.35), Inches(11.5), Inches(0.95))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    
    # Category Tag (Muted Yellow, Small Caps)
    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Trebuchet MS"
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_YELLOW
    p_cat.space_after = Pt(2)
    
    # Slide Title (Lime Green, All Caps, Bold)
    p_title = tf.add_paragraph()
    p_title.text = title_text.upper()
    p_title.font.name = "Trebuchet MS"
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_LIME

def add_card(slide, left, top, width, height, card_title, bullets, title_color=COLOR_YELLOW):
    """Draws a premium dark card container with a neon green outline and structured text inside."""
    # Card Background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_LIME
    card.line.width = Pt(1.2)
    
    # Content Textbox (adding padding)
    tx_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    # Card Title
    p_title = tf.paragraphs[0]
    p_title.text = card_title.upper()
    p_title.font.name = "Trebuchet MS"
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = title_color
    p_title.space_after = Pt(12)
    
    # Bullets
    for i, bullet in enumerate(bullets):
        p_bullet = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        # If it's the first paragraph, keep card title space_after, else add normal paragraph
        if i > 0:
            p_bullet = tf.add_paragraph()
        
        # Check for bold prefixes separated by ":"
        if ":" in bullet:
            parts = bullet.split(":", 1)
            p_bullet.text = ""
            run_bold = p_bullet.add_run()
            run_bold.text = parts[0] + ":"
            run_bold.font.name = "Calibri"
            run_bold.font.size = Pt(13)
            run_bold.font.bold = True
            run_bold.font.color.rgb = COLOR_WHITE
            
            run_norm = p_bullet.add_run()
            run_norm.text = parts[1]
            run_norm.font.name = "Calibri"
            run_norm.font.size = Pt(13)
            run_norm.font.color.rgb = COLOR_MUTED
        else:
            p_bullet.text = bullet
            p_bullet.font.name = "Calibri"
            p_bullet.font.size = Pt(13)
            p_bullet.font.color.rgb = COLOR_MUTED
            
        p_bullet.space_after = Pt(8)

# ----------------- SLIDE 1: TITLE SLIDE -----------------
slide_layout = prs.slide_layouts[6]  # Blank layout
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)

# Large Center Title Card
title_left = Inches(1.5)
title_top = Inches(1.8)
title_width = Inches(10.33)
title_height = Inches(4.0)

# Neon green frame for title slide
title_frame = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, title_left, title_top, title_width, title_height)
title_frame.fill.solid()
title_frame.fill.fore_color.rgb = COLOR_CARD_BG
title_frame.line.color.rgb = COLOR_LIME
title_frame.line.width = Pt(3)

# Title Text frame
tx_title = slide1.shapes.add_textbox(title_left + Inches(0.4), title_top + Inches(0.4), title_width - Inches(0.8), title_height - Inches(0.8))
tf1 = tx_title.text_frame
tf1.word_wrap = True

p_main = tf1.paragraphs[0]
p_main.text = "FEMFIT - INDUSTRIAL"
p_main.alignment = PP_ALIGN.CENTER
p_main.font.name = "Trebuchet MS"
p_main.font.size = Pt(54)
p_main.font.bold = True
p_main.font.color.rgb = COLOR_LIME
p_main.space_after = Pt(14)

p_sub = tf1.add_paragraph()
p_sub.text = "Edge-AI Smart Wearable Vest Calibrated for Female Ergonomics"
p_sub.alignment = PP_ALIGN.CENTER
p_sub.font.name = "Calibri"
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = COLOR_WHITE
p_sub.space_after = Pt(28)

p_team = tf1.add_paragraph()
p_team.text = "DEVELOPED BY: THE GILDED GIRL  |  KOMAL HARSHITA"
p_team.alignment = PP_ALIGN.CENTER
p_team.font.name = "Trebuchet MS"
p_team.font.size = Pt(14)
p_team.font.bold = True
p_team.font.color.rgb = COLOR_YELLOW

# ----------------- SLIDE 2: PROBLEM STATEMENT -----------------
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
create_slide_header(slide2, "The Ergonomic Safety Gap")

add_card(slide2, 
         left=Inches(0.6), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="The Core Problem",
         bullets=[
             "Widespread MSD Injury: Musculoskeletal Disorders (MSDs) are the leading cause of industrial worker absenteeism globally.",
             "The Anthropometric Bias: Standard safety vests and RULA/REBA guidelines assume male center of gravity and pelvic anatomy.",
             "Female Lifting Vulnerability: Females have different pelvic loading limits, meaning standard 'lifts' still cause high spine shear forces.",
             "Mechanical Vibrations: Continuous micro-vibrations from machinery damage vascular and nerve networks over time."
         ])

add_card(slide2, 
         left=Inches(6.9), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="Market Size & Impact",
         bullets=[
             "Target User Base: Over 150 Million female industrial and warehouse workers globally.",
             "Financial Strain: MSDs cost employers $50 Billion+ annually in worker compensation and lost productivity.",
             "Smart Wearable Market: Industrial IoT safety wearable market projected to cross $5 Billion by 2028.",
             "Hackathon Objective: Deliver a functional, low-cost ($20 unit cost) Edge-AI vest to prevent injury at the source."
         ], title_color=COLOR_LIME)

# ----------------- SLIDE 3: OBJECTIVE & APPROACH -----------------
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
create_slide_header(slide3, "Objective & Methodology")

# Left box - Objective
add_card(slide3,
         left=Inches(0.6), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="Project Objectives",
         bullets=[
             "Edge AI Detection: Implement zero-latency local classification of physical activities and safety risks.",
             "Immediate Bio-Feedback: Warn workers within 50ms of detecting an anomaly to immediately correct posture.",
             "Continuous Machinery Assessment: Track high-frequency vibrations to warn operators of vascular damage risks.",
             "Edge Autonomy: Eliminate dependency on cloud servers, ensuring operation in remote, offline industrial environments."
         ], title_color=COLOR_LIME)

# Right box - Methodology
add_card(slide3,
         left=Inches(6.9), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="Technical Approach",
         bullets=[
             "1. Sensor Fusion: Read 6-axis IMU data (50Hz) mounted on the spine to get acceleration and angular velocities.",
             "2. Local Processing: Estimate precise pitch (trunk flexion) and roll (asymmetric twist) using a drift-free complementary filter.",
             "3. Vibration Analytics: Calculate rolling window variance to isolate machine-induced micro-vibrations from body movements.",
             "4. Ergonomic Decision Matrix: Apply RULA/REBA thresholds calibrated to female lumbar limits to command haptic vibrations."
         ])

# ----------------- SLIDE 4: SOLUTION OVERVIEW (1/2) -----------------
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4)
create_slide_header(slide4, "FemFit-Industrial Solution Overview", "Solution Architecture")

# 3 Columns of features
add_card(slide4,
         left=Inches(0.6), top=Inches(1.6), width=Inches(3.8), height=Inches(5.2),
         card_title="Hardware Components",
         bullets=[
             "Microcontroller: ESP32 or Nano 33 BLE Sense processing local streams.",
             "6-Axis IMU: Onboard MPU6050 tracking linear and angular kinetics.",
             "Haptic Motor: High-amplitude ERM vibration disk for instant feedback.",
             "Power: Ultra-lightweight LiPo battery with power-saving deep sleep integration."
         ])

add_card(slide4,
         left=Inches(4.75), top=Inches(1.6), width=Inches(3.8), height=Inches(5.2),
         card_title="Zero-Cloud Edge AI",
         bullets=[
             "Heuristic Phase: Calibrated geometric decision matrix acting as the PoC classifier.",
             "TinyML Upgrade: TensorFlow Lite for Microcontrollers (TFLite Micro) pipeline.",
             "Ultra-Low Latency: On-chip processing executes calculations in under 5 milliseconds."
         ], title_color=COLOR_LIME)

add_card(slide4,
         left=Inches(8.9), top=Inches(1.6), width=Inches(3.8), height=Inches(5.2),
         card_title="Ergonomic Customization",
         bullets=[
             "Gender-specific tuning: Tailored pelvic limits and forward bending angles.",
             "Smart Warning: Distinct vibration alerts separating static posture limits from machine jitter.",
             "Form-Factor: Integrated directly into standard mesh safety vests."
         ])

# ----------------- SLIDE 5: SOLUTIONS - BIOMECHANICS -----------------
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5)
create_slide_header(slide5, "Female Ergonomic Calibration", "Biomechanical Rationale")

add_card(slide5,
         left=Inches(0.6), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="RULA/REBA Adjustments",
         bullets=[
             "Lower Center of Gravity: Females balance load differently, making deep waist bends (flexion) significantly more stressful on lumbar discs (L4-S1).",
             "Trunk Flexion (Pitch) Limits: Warning at >20 degrees (REBA score 2). Critical alert at >45 degrees (REBA score 3/4) to prevent spinal shear.",
             "Spinal Twisting (Roll/Yaw): Warning at >15 degrees. Combined bending and twisting is the #1 cause of acute disc herniations."
         ])

add_card(slide5,
         left=Inches(6.9), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="Vibration Strain Limits",
         bullets=[
             "Vibration Hazards: Heavy machinery exposes workers to micro-vibrations between 10 Hz and 50 Hz.",
             "Vascular Strain: ISO 5349 limits define sustained exposure as a trigger for Hand-Arm Vibration Syndrome.",
             "Detection Logic: High-pass filtering isolated from postural sway. If high-frequency variance exceeds 0.3g RMS, the system triggers the alert."
         ], title_color=COLOR_LIME)

# ----------------- SLIDE 6: TECHNICAL IMPLEMENTATION -----------------
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6)
create_slide_header(slide6, "Technical Implementation Stack")

add_card(slide6,
         left=Inches(0.6), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="Data & Model Pipeline",
         bullets=[
             "Simulation: High-fidelity Python synthetic data generator recreating 8-hour shifts.",
             "Edge Impulse: CSV uploads structured in class folders for rapid spectral processing.",
             "DSP Spectral Block: Filters noise and extracts frequency power spectrum from IMU signals.",
             "TFLite Classifier: Neural network trained in Edge Impulse and exported as a lightweight C++ library."
         ], title_color=COLOR_LIME)

add_card(slide6,
         left=Inches(6.9), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="On-Device C++ Architecture",
         bullets=[
             "IDE: Built for Arduino IDE / ESP-IDF compiler.",
             "Drift Mitigation: Complementary Filter combining accelerometer gravity vectors and gyroscope integral rates.",
             "Non-Blocking Scheduler: Uses task loops (simulating RTOS) to query sensors and pulse haptics simultaneously without delays.",
             "Safety Logic: Real-time decision matrix evaluates thresholds instantly."
         ])

# ----------------- SLIDE 7: CHALLENGES & MITIGATIONS -----------------
slide7 = prs.slides.add_slide(slide_layout)
set_slide_background(slide7)
create_slide_header(slide7, "Challenges Faced & Mitigations")

# 3 Horizontal Cards (using left, top, width, height)
add_card(slide7,
         left=Inches(0.6), top=Inches(1.6), width=Inches(3.8), height=Inches(5.2),
         card_title="1. Sensor Drift",
         bullets=[
             "Challenge: Gyroscope integration drifts over time, resulting in false posture warnings after minutes.",
             "Mitigation: Implemented a fast-converging complementary filter ($0.98$ gyro + $0.02$ accel) to bound pitch/roll errors to 1 degree."
         ])

add_card(slide7,
         left=Inches(4.75), top=Inches(1.6), width=Inches(3.8), height=Inches(5.2),
         card_title="2. Latency Limitations",
         bullets=[
             "Challenge: If alert processing takes >100ms, the worker has completed the lift, rendering the warning useless.",
             "Mitigation: Replaced heavy floating-point algebra with fixed-point decision matrix, executing calculations in <5ms."
         ], title_color=COLOR_LIME)

add_card(slide7,
         left=Inches(8.9), top=Inches(1.6), width=Inches(3.8), height=Inches(5.2),
         card_title="3. Data Scarcity",
         bullets=[
             "Challenge: Gathering real IMU training data for female posture anomalies takes months.",
             "Mitigation: Built a synthetic simulation script in Python to generate thousands of labeled trials representing clean baselines and specific strains."
         ])

# ----------------- SLIDE 8: RESULTS & ACHIEVEMENTS -----------------
slide8 = prs.slides.add_slide(slide_layout)
set_slide_background(slide8)
create_slide_header(slide8, "Outcomes and Achievements")

add_card(slide8,
         left=Inches(0.6), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="Expected Impact",
         bullets=[
             "Injury Reduction: Targeted 60%+ drop in back muscle strain and fatigue-related MSD cases over 12 months.",
             "Immediate Correction: Over 90% compliance rate expected as workers instinctively correct posture during haptic pulsing.",
             "Operational Efficiency: Lower absenteeism, leading to an estimated 15% increase in shift productivity.",
             "ROI: Vest pays for itself in less than 3 months by avoiding a single MSD workers' comp claim."
         ], title_color=COLOR_LIME)

add_card(slide8,
         left=Inches(6.9), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="Key Accomplishments",
         bullets=[
             "High-fidelity Simulator: Successfully generated 120 dataset trials mapping diverse ergonomic movements.",
             "Robust Firmware: Created a dual-target Arduino firmware (ESP32 / BLE Sense) with complementary filter.",
             "Modular Code: Handled non-blocking haptic state logic for distinct warning patterns.",
             "TinyML Ready: Established a seamless transition path to Edge Impulse."
         ])

# ----------------- SLIDE 9: DEMONSTRATION & LOGIC -----------------
slide9 = prs.slides.add_slide(slide_layout)
set_slide_background(slide9)
create_slide_header(slide9, "System Prototype & Verification")

add_card(slide9,
         left=Inches(0.6), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="Haptic Pulse Grammar",
         bullets=[
             "Visual Feedback: Microcontroller continuously reads IMU data and processes sensor orientation.",
             "Warning 1: Posture Fix (Pitch 20-45 deg) -> Moderate short pulses (200ms ON, 400ms OFF) prompting correction.",
             "Warning 2: Critical Strain (Pitch >45 deg OR Roll >20 deg) -> Rapid warning double pulses (100ms ON/OFF/ON, 500ms OFF).",
             "Warning 3: Machinery Jitter -> Extended sustained vibrations (800ms ON, 200ms OFF) warning of vibration overload."
         ])

add_card(slide9,
         left=Inches(6.9), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="Filter & Math Verification",
         bullets=[
             "Complementary Filter: $Pitch = 0.98*(Pitch + Gyro*dt) + 0.02*AccelAngle$ guarantees zero gyro drift while filtering high-frequency noise.",
             "Vibration RMS: Sliding window variance $Var = \\Sigma(a_i - \\mu)^2 / N$ captures mechanical vibrations from 10Hz to 25Hz.",
             "Execution Time: Under 4ms per loop iteration, easily beating the <50ms hackathon latency budget."
         ], title_color=COLOR_LIME)

# ----------------- SLIDE 10: FUTURE ENHANCEMENTS & PLAN -----------------
slide10 = prs.slides.add_slide(slide_layout)
set_slide_background(slide10)
create_slide_header(slide10, "Project Roadmap & Future Plan")

add_card(slide10,
         left=Inches(0.6), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="Future Enhancements",
         bullets=[
             "1. Edge Impulse Deploy: Train a 1D Convolutional Neural Network (CNN) on Edge Impulse, compile it as a TFLite Micro model, and deploy it to replace the heuristic C++ code.",
             "2. Custom Hardware: Shrink ESP32 system into a custom coin-sized PCB integrated with ultra-thin flexible piezo actuators.",
             "3. BLE Telemetry: Stream real-time anomaly telemetry to an industrial dashboard for safety supervisor review.",
             "4. Multi-Sensor Mesh: Place secondary IMU node on the worker's thigh to verify squat vs. stoop biomechanics."
         ])

add_card(slide10,
         left=Inches(6.9), top=Inches(1.6), width=Inches(5.8), height=Inches(5.2),
         card_title="Prototype Project Plan",
         bullets=[
             "Milestone 1 (Hackathon Day 1): Complete data generation, write C++ sensor drivers, and verify pitch/roll complementary filter.",
             "Milestone 2 (Hackathon Day 2): Tune RULA/REBA threshold matrices and test the haptic alarm state machine on hardware.",
             "Milestone 3 (Post-Hackathon Wk 1): Ingest synthetic dataset into Edge Impulse and deploy the compiled TFLite Micro model.",
             "Milestone 4 (Post-Hackathon Wk 2-3): Conduct field tests with 10 female warehouse operators to refine biomechanical angles."
         ], title_color=COLOR_LIME)

# 4. Save Presentation
output_filename = "femfit_presentation.pptx"
output_path = os.path.join(os.getcwd(), output_filename)
prs.save(output_path)
print(f"Presentation saved successfully at: {output_path}")
